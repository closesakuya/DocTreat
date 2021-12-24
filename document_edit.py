#!/usr/bin/env python
# -*- coding:utf-8 -*-
import os
import time

import win32com
from win32com.client import Dispatch
from win32com import client as wc
import pythoncom
import openpyxl
import re
# import pptx
import docx
from pptx import Presentation
from typing import List, Dict, Tuple
import datetime


class EditDocx:
    def __init__(self, filename=None):
        self.doc = None
        self.filename = None
        if filename:
            self.filename = filename
            if os.path.exists(self.filename):
                self.doc = docx.Document(filename)
            else:
                raise NotImplemented
        else:
            raise NotImplemented

    def __del__(self):
        if self.doc:
            pass

    @staticmethod
    def _replace_text(text_frame, replace_dct: Dict, replace_cb_lst: List = None):
        cnt = 0
        for paragraph in text_frame.paragraphs:
            val = paragraph.text
            if replace_cb_lst:
                for func in replace_cb_lst:
                    val, is_change = func(val)
                    if is_change:
                        paragraph.text = val
                        cnt += 1
            else:
                for string, new_string in replace_dct.items():
                    res = re.search(string, val)
                    if res:
                        val = re.sub(string, new_string, val)
                        paragraph.text = val
                        cnt += 1

        return cnt

    @staticmethod
    def _replace_table(table, replace_dct: Dict, replace_cb_lst: List = None):
        cnt = 0
        for row in range(table.rows.__len__()):
            for cell in table.rows[row].cells:
                val = cell.text
                print(val, replace_dct)
                if replace_cb_lst:
                    for func in replace_cb_lst:
                        val, is_change = func(val)
                        if is_change:
                            cell.text = val
                            cnt += 1
                else:
                    for string, new_string in replace_dct.items():
                        print(string, val)
                        res = re.search(string, val)
                        if res:
                            val = re.sub(string, new_string, val)
                            cell.text = val
                            cnt += 1

        return cnt

    def replace_str(self, replace_dct: Dict):
        """替换文字"""
        cnt = 0
        for para in self.doc.paragraphs:
            for i in range(len(para.runs)):
                val = para.runs[i].text
                for string, new_string in replace_dct.items():
                    if string in val:
                        val = val.replace(string, new_string)
                        para.runs[i].text = val
                        cnt += 1
        # 页眉页脚
        for i in range(self.doc.sections.__len__()):
            header = self.doc.sections[i].header
            footer = self.doc.sections[i].footer
            for para in header.paragraphs:
                for j in range(len(para.runs)):
                    val = para.runs[j].text
                    for string, new_string in replace_dct.items():
                        if string in val:
                            val.replace(string, new_string)
                            para.runs[j].text = val
                            cnt += 1
            for para in footer.paragraphs:
                for j in range(len(para.runs)):
                    val = para.runs[j].text
                    for string, new_string in replace_dct.items():
                        if string in val:
                            val.replace(string, new_string)
                            para.runs[j].text = val
                            cnt += 1
        return cnt

    def replace(self, replace_dct: Dict, replace_cb_lst: List = None):
        """采用通配符匹配替换"""
        cnt = 0
        cnt += self._replace_text(self.doc, replace_dct, replace_cb_lst)
        for i in range(self.doc.tables.__len__()):
            cnt += self._replace_table(self.doc.tables[i], replace_dct, replace_cb_lst)

        # for para in self.doc.paragraphs:
        #     val = para.text
        #     if replace_cb_lst:
        #         for func in replace_cb_lst:
        #             val, is_change = func(val)
        #             if is_change:
        #                 para.text = val
        #                 cnt += 1
        #     else:
        #         for string, new_string in replace_dct.items():
        #             res = re.search(string, val)
        #             if res:
        #                 val = re.sub(string, new_string, val)
        #                 cnt += 1
        #             para.text = val


        # 页眉页脚
        for i in range(self.doc.sections.__len__()):
            header = self.doc.sections[i].header
            footer = self.doc.sections[i].footer
            cnt += self._replace_text(header, replace_dct, replace_cb_lst)
            cnt += self._replace_text(footer, replace_dct, replace_cb_lst)
            for i in range(header.tables.__len__()):
                cnt += self._replace_table(header.tables[i], replace_dct, replace_cb_lst)
            for i in range(footer.tables.__len__()):
                cnt += self._replace_table(footer.tables[i], replace_dct, replace_cb_lst)
            # for para in header.paragraphs:
            #     val = para.text
            #     print(val)
            #     for j in range(len(para.runs)):
            #         val2 = para.runs[j].text
            #         print(val2)
            #     if replace_cb_lst:
            #         for func in replace_cb_lst:
            #             val, is_change = func(val)
            #             if is_change:
            #                 para.text = val
            #                 cnt += 1
            #     else:
            #         for string, new_string in replace_dct.items():
            #             res = re.search(string, val)
            #             if res:
            #                 val = re.sub(string, new_string, val)
            #                 para.text = val
            #                 cnt += 1
            #
            # for para in footer.paragraphs:
            #     val = para.text
            #     if replace_cb_lst:
            #         for func in replace_cb_lst:
            #             val, is_change = func(val)
            #             if is_change:
            #                 para.text = val
            #                 cnt += 1
            #     else:
            #         for string, new_string in replace_dct.items():
            #             res = re.search(string, val)
            #             if res:
            #                 val = re.sub(string, new_string, val)
            #                 para.text = val
            #                 cnt += 1

        return cnt

    def save(self):
        """保存文档"""
        self.doc.save(self.filename)

    def save_as(self, filename):
        """文档另存为"""
        self.doc.save(filename)


class EditXlsx:
    def __init__(self, filename=None):
        self.wb = None
        self.filename = None
        if filename:
            self.filename = filename
            if os.path.exists(self.filename):
                self.wb = openpyxl.load_workbook(filename, data_only=True)
            else:
                raise NotImplemented
        else:
            raise NotImplemented

    def __del__(self):
        if self.wb:
            self.wb.close()

    def replace_str(self, replace_dct: Dict):
        """替换文字"""
        cnt = 0
        for item in self.wb.sheetnames:
            cur = self.wb[item]
            for i, row in enumerate(cur.iter_rows()):
                for j, cell in enumerate(row):
                    val = str(cell.value)
                    for string, new_string in replace_dct.items():
                        if string in val:
                            val = val.replace(string, new_string)
                            cell.value = val
                            cnt += 1
        return cnt

    def replace(self, replace_dct: Dict, replace_cb_lst: List = None):
        """采用通配符匹配替换"""
        cnt = 0
        for item in self.wb.sheetnames:
            cur = self.wb[item]
            for i, row in enumerate(cur.iter_rows()):
                for j, cell in enumerate(row):
                    val = str(cell.value)
                    if replace_cb_lst:
                        for func in replace_cb_lst:
                            val, is_change = func(val)
                            if is_change:
                                cell.value = val
                                cnt += 1
                    else:
                        for string, new_string in replace_dct.items():
                            res = re.search(string, val)
                            if res:
                                val = re.sub(string, new_string, val)
                                cell.value = val
                                cnt += 1
        return cnt

    def save(self):
        """保存文档"""
        self.wb.save(self.filename)

    def save_as(self, filename):
        """文档另存为"""
        self.wb.save(filename)


class EditPptx:
    def __init__(self, filename=None):
        self.prs = None
        self.filename = None
        if filename:
            self.filename = filename
            if os.path.exists(self.filename):
                self.prs = Presentation(filename)
            else:
                raise NotImplemented
        else:
            raise NotImplemented

    def __del__(self):
        if self.prs:
            pass

    @staticmethod
    def _replace_text(text_frame, replace_dct: Dict, replace_cb_lst: List = None):
        cnt = 0
        for paragraph in text_frame.paragraphs:
            val = paragraph.text
            if replace_cb_lst:
                for func in replace_cb_lst:
                    val, is_change = func(val)
                    if is_change:
                        paragraph.text = val
                        cnt += 1
            else:
                for string, new_string in replace_dct.items():
                    res = re.search(string, val)
                    if res:
                        val = re.sub(string, new_string, val)
                        paragraph.text = val
                        cnt += 1

        return cnt

    def replace_str(self, replace_dct: Dict):
        """替换文字"""
        return self.replace(replace_dct)

    def replace(self, replace_dct, replace_cb_lst: List = None):
        """采用通配符匹配替换"""
        cnt = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:  # 判断Shape是否含有文本框
                    text_frame = shape.text_frame
                    cnt += self._replace_text(text_frame, replace_dct, replace_cb_lst)  # 调用replace_text函数实现文本替换
                if shape.has_table:  # 判断Shape是否含有表格
                    table = shape.table
                    for cell in table.iter_cells():  # 遍历表格的cell
                        text_frame = cell.text_frame
                        cnt += self._replace_text(text_frame, replace_dct, replace_cb_lst)  # 调用replace_text函数实现文本替换
        return cnt

    def save(self):
        """保存文档"""
        self.prs.save(self.filename)

    def save_as(self, filename):
        """文档另存为"""
        self.prs.save(filename)


class EditText:
    def __init__(self, filename=None):
        self.filename = None
        self.content = []
        self.my_encoding = 'utf-8'
        if filename:
            self.filename = filename
            try:
                # 自动识别文件编码
                import chardet
                with open(self.filename, 'rb') as f:
                    self.my_encoding = chardet.detect(f.read(200))['encoding']
                with open(filename, "r+", encoding=self.my_encoding) as f:
                    self.content = f.readlines()
            except Exception as e:
                print(e)
                self.content = []
        else:
            raise NotImplemented

    def __del__(self):
        pass

    def replace_str(self, replace_dct: Dict):
        """替换文字"""
        cnt = 0
        if not self.content:
            return 0
        for i in range(self.content.__len__()):
            txt = self.content[i]
            for string, new_string in replace_dct.items():
                if string in txt:
                    txt = txt.replace(str, new_string)
                    self.content[i] = txt
                    cnt += 1
        return cnt

    def replace(self, replace_dct: Dict, replace_cb_lst: List = None):
        """采用通配符匹配替换"""
        cnt = 0
        if not self.content:
            return 0
        for i in range(self.content.__len__()):
            val = self.content[i]
            if replace_cb_lst:
                for func in replace_cb_lst:
                    val, is_change = func(val)
                    if is_change:
                        self.content[i] = val
                        cnt += 1
            else:
                for string, new_string in replace_dct.items():
                    res = re.search(string, val)
                    if res:
                        val = re.sub(string, new_string, val)
                        self.content[i] = val
                        cnt += 1
        return cnt

    def save(self):
        """保存文档"""
        if self.content:
            try:
                with open(self.filename, "w+", encoding=self.my_encoding) as f:
                    f.writelines(self.content)
            except Exception as e:
                print(e)

    def save_as(self, filename):
        """文档另存为"""
        if self.content:
            try:
                with open(filename, "w+", encoding=self.my_encoding) as f:
                    f.writelines(self.content)
            except Exception as e:
                print(e)


def get_edit(file_name: str):
    sub_name = os.path.splitext(file_name)[-1].lower()
    m_map = {
        ".docx": EditDocx,
        ".doc": EditDocx,
        ".xlsx": EditXlsx,
        ".xls": EditXlsx,
        ".pptx": EditPptx,
        ".ppt": EditPptx,
        ".csv": EditText,
        ".txt": EditText
    }
    return m_map.get(sub_name, EditText)


_y_map = {"0": "零", "1": "一", "2": "二", "3": "三", "4": "四",
         "5": "五", "6": "六", "7": "七", "8": "八", "9": "九"}
_y_map_2 = {"0": "〇", "1": "一", "2": "二", "3": "三", "4": "四",
         "5": "五", "6": "六", "7": "七", "8": "八", "9": "九"}
_m_map = {"1": "一", "2": "二", "3": "三", "4": "四",
         "5": "五", "6": "六", "7": "七", "8": "八", "9": "九",
         "01": "一", "02": "二", "03": "三", "04": "四",
         "05": "五", "06": "六", "07": "七", "08": "八", "09": "九",
         "10": "十", "11": "十一", "12": "十二"}

_m_simple_month_days_count = [31, 29, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31]


def replace_date(input_str: str, ori_date: List, dst_date: List,
                 fmt_index_used: Tuple[int] = (0, 1), for_excel_sp: bool = False):
    """
    TODO 如果excel中输入‘2021年1月’ 则调用openexls获取的value为1900年1月1日开始的日期数int,
    暂时判断int值是否符合，可能会替换掉正好为该值的int类型值
    """
    # print("func ", ori_date, dst_date, fmt_index_used, input_str, "aaa")
    assert ori_date.__len__() == dst_date.__len__()
    ori_date_str = [str(ori_date[0]), str(ori_date[1])]
    dst_date_str = [str(dst_date[0]), str(dst_date[1])]
    y_str = "".join(_y_map[item] for item in str(ori_date[0]))
    y_str_2 = "".join(_y_map_2[item] for item in str(ori_date[0]))
    dst_m_str = _m_map[str(dst_date[1])]
    ret_str = ""

    # 特殊处理,excel中表格日期
    if for_excel_sp and input_str.isdigit():
        src_stamp_day = int(input_str)
        cmp_stamp_day = (datetime.date(ori_date[0], ori_date[1], 1) - datetime.date(1900, 1, 1)).days + 2
        time_dlt = src_stamp_day - cmp_stamp_day
        # print("date is int ", cmp_stamp_day, int(input_str), ori_date[0], ori_date[1])
        if 0 <= time_dlt <= _m_simple_month_days_count[ori_date[1] - 1]:
            input_day = datetime.date(ori_date[0], ori_date[1], 1) + datetime.timedelta(days=time_dlt)
            ret_str = "{0}年{1}月{2}日".format(dst_date[0], dst_date[1], input_day.day)
            return ret_str, True

    if ori_date.__len__() == 2:
        mc_fmts = [
            r"(" + str(ori_date[0]) + r")" +
            r"(\s?[\.\-\s\\/年]{1,1}\s?)(" + str(ori_date[1]) + r"|" + "{0:02d}".format(ori_date[1]) +
            r")(\s?[\s\.\-\\/月]{0,1}[0-9]{0,2})",
            r"(" + y_str + r"|" + y_str_2 + r")" +
            r"(\s?[\s\.年]{1,1}\s?)(" + _m_map[str(ori_date[1])] + r")(\s?[\s\.月$\r\n]{1,1})"]
        mc_fmts = [mc_fmts[i] for i in fmt_index_used]
        found = False
        ret_str = input_str
        # if input_str and input_str != "None":
        #     print(mc_fmts, input_str)
        for fmt in mc_fmts:
            if re.search(fmt, ret_str):
                # print("FFFF ", re.split(fmt, ret_str))
                found = True
                words = re.split(fmt, ret_str)
                i = 0
                # print(words)
                while i < (words.__len__() - 3):
                    if not words[i]:
                        i += 1
                        continue
                    per_word = "".join([words[i + per] for per in range(4)])
                    if re.match(fmt, per_word):
                        # print("aaa", per_word, ori_date_str, i, words)
                        if words[i + 0] == ori_date_str[0]:  # 2021
                            if words[i + 2].__len__() == 2:  # 2021.02
                                words[i + 0] = dst_date_str[0]
                                words[i + 2] = "{0:02d}".format(dst_date[1])
                            else:  # 2021.02
                                words[i + 0] = dst_date_str[0]
                                words[i + 2] = "{0:d}".format(dst_date[1])
                        else:  # 二零二一
                            if re.search(r'〇', words[i + 0]):  # 二〇二一
                                words[i + 0] = "".join(_y_map_2[item] for item in dst_date_str[0])
                                words[i + 2] = dst_m_str
                            else:
                                words[i + 0] = "".join(_y_map[item] for item in dst_date_str[0])
                                words[i + 2] = dst_m_str
                        i += 4
                    else:
                        i += 1
                # print(words)
                ret_str = "".join(words)
                # print("Done", ret_str)
    else:
        raise NotImplementedError  # TODO
    return ret_str, found


def replace_dir(dir_path: str, replace_dct: Dict, filter_lst: List = None, logger=None, replace_cb_lst: List = None):
    if os.path.isdir(dir_path):
        tt_cnt = 0
        for item in os.listdir(dir_path):
            abs_path = os.path.join(dir_path, item)
            tt_cnt += replace_dir(abs_path, replace_dct, filter_lst, logger=logger, replace_cb_lst=replace_cb_lst)

        return tt_cnt
    elif os.path.isfile(dir_path):
        filename = os.path.split(dir_path)[-1]
        found = False
        if filter_lst:
            for item in filter_lst:
                if re.search(item, filename):
                    found = True
                    break
        if found:
            if re.search(r"~.*\..*", filename):  # 去除~开头的文件
                found = False
            else:
                found = True
        if found:
            logger("开始替换文本: 【{0}】".format(filename))
            cnt = 0
            try:
                inst = get_edit(file_name=filename)(dir_path)
                cnt = inst.replace(replace_dct, replace_cb_lst=replace_cb_lst)
                inst.save()
                logger("替换关键字次数【{0}】, 文件【{1}】".format(cnt, filename))
            except Exception as e1:
                logger("e: {0} file: {1}".format(e1, os.path.split(dir_path)[-1]))
                raise e1
            return cnt
    return 0


def office2officeX(path, logger=print):
    path = path.replace("/", os.sep).replace("\\", os.sep)
    for subPath in os.listdir(path):
        subPath = os.path.join(path, subPath)
        if os.path.isdir(subPath):
            office2officeX(subPath, logger=logger)
        elif re.search(r"~.*\..*", subPath):  # 去除~开头的文件
            continue
        elif subPath.endswith('.ppt'):
            if os.path.exists("{}x".format(subPath)):
                continue
            pythoncom.CoInitialize()
            logger("正在转换文件{0} to pptx...".format(os.path.split(subPath)[-1]))
            powerpoint = win32com.client.Dispatch('PowerPoint.Application')
            win32com.client.gencache.EnsureDispatch('PowerPoint.Application')
            powerpoint.Visible = 1
            ppt = powerpoint.Presentations.Open(subPath)
            ppt.SaveAs("{}x".format(subPath))
            ppt.Close()
            try:
                powerpoint.Quit()  # 启动报错
                os.remove(subPath)
            except Exception as e:
                raise e
        elif subPath.endswith('.doc'):
            if os.path.exists("{}x".format(subPath)):
                continue
            pythoncom.CoInitialize()
            logger("正在转换文件{0} to docx...".format(os.path.split(subPath)[-1]))
            time.sleep(1)  # 此处不sleep连续打开word会报错 TODO
            word = wc.Dispatch("Word.Application")  # 打开word应用程序
            doc = word.Documents.Open(subPath)  # 打开word文件
            doc.SaveAs("{}x".format(subPath), FileFormat=12)  # 另存为后缀为".docx"的文件，其中参数12指docx文件
            doc.Close()  # 关闭原来word文件
            try:
                word.Quit()
                os.remove(subPath)
            except Exception as e:
                raise e
        elif subPath.endswith('.xls'):
            if os.path.exists("{}x".format(subPath)):
                continue
            pythoncom.CoInitialize()
            logger("正在转换文件{0} to xlsx...".format(os.path.split(subPath)[-1]))
            excel = wc.Dispatch("Excel.Application")  # 打开word应用程序
            wb = excel.Workbooks.Open(subPath)
            wb.SaveAs("{}x".format(subPath), FileFormat=51)  # 另存为后缀为".docx"的文件，其中参数12指docx文件
            wb.Close()
            try:
                excel.Quit()
                os.remove(subPath)
            except Exception as e:
                raise e


# if __name__ == "__main__":
    # office2officeX(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\YYY')
    # doc2docx(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\XXX\aa.doc')
    # e = EditXlsx(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\YYY\AAA.xlsx')
    # e = EditDocx(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\YYY\AAA.docx')
    # e = EditPptx(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\YYY\AAA.pptx')
    # e = EditText(r'D:\share_dir\product_env\01.SVN\01.local_git\DocTreat\YYY\AAA.csv')
    # # e.replace_str("HAHA", "sss")
    # e.replace({r'([0-9].*)UUU[A-Za-z].*': 'HELLO', 'D': 'HELLO'})
    # e.save()

    # powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    # powerpoint.Visible = 1
    # subPath = "D:/share_dir/product_env/01.SVN/01.local_git/DocTreat/YYY/AAA - 副本 (2).ppt"
    # subPath = subPath.replace("/",os.sep)
    # print(subPath)
    # ppt = powerpoint.Presentations.Open(subPath)

# msg, is_change = replace_date("\n2021年11月11日\n", (2021, 11), (2021, 8))
# print(msg, is_change)

